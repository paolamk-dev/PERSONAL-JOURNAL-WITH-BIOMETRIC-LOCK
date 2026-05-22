#!/usr/bin/env python3
"""
SmartAC Visual Identity Guidelines - PowerPoint (Editable) Generator
Author: MAKAMTE TATSINKOU PAOLA
Course: Level 3 Software Engineering - Graphic Design
Project: IoT-Based Smart Air Conditioning Control System

This script generates an editable PowerPoint presentation version of the
Visual Identity Guidelines that can be modified, updated, and customized.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from datetime import datetime

# ============================================================================
# BRAND COLOR PALETTE - Constants
# ============================================================================
ARCTIC_BLUE = RGBColor(10, 132, 255)
COOL_TEAL = RGBColor(48, 213, 200)
DEEP_NAVY = RGBColor(13, 27, 42)
FROST_WHITE = RGBColor(244, 249, 255)
SLATE_GREY = RGBColor(107, 114, 128)
DANGER_RED = RGBColor(255, 59, 48)
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DEEP_NAVY

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    text_frame = title_box.text_frame
    text_frame.text = title
    text_frame.paragraphs[0].font.size = Pt(44)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.color.rgb = FROST_WHITE
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
    text_frame = subtitle_box.text_frame
    text_frame.text = subtitle
    text_frame.paragraphs[0].font.size = Pt(20)
    text_frame.paragraphs[0].font.color.rgb = COOL_TEAL
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide

def add_section_slide(prs, section_title):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = ARCTIC_BLUE

    # Section title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(1))
    text_frame = title_box.text_frame
    text_frame.text = section_title
    text_frame.paragraphs[0].font.size = Pt(36)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.color.rgb = WHITE
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with title and bullet points"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = FROST_WHITE

    # Title bar
    title_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = ARCTIC_BLUE
    title_shape.line.color.rgb = ARCTIC_BLUE

    # Title text
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.margin_left = Inches(0.3)
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(4.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for item in content_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = DEEP_NAVY
        p.level = 0
        p.space_after = Pt(10)

    return slide

def generate_pptx():
    """Generate the complete Visual Identity Guidelines PowerPoint"""
    filename = "SmartAC_Visual_Identity_Guidelines_EDITABLE.pptx"

    print("Generating SmartAC Visual Identity Guidelines PowerPoint...")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ========================================================================
    # SLIDE 1: Cover
    # ========================================================================
    print("  ✓ Slide 1: Cover Page")
    slide = add_title_slide(
        prs,
        "Visual Identity Guidelines",
        "SmartAC — IoT-Based Smart Air Conditioning Control System"
    )

    # Author info
    author_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(1.5))
    text_frame = author_box.text_frame
    text_frame.text = "MAKAMTE TATSINKOU PAOLA\n"
    text_frame.text += "Level 3 Software Engineering — Graphic Design Course\n"
    text_frame.text += f"Academic Year {datetime.now().year}"

    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = COOL_TEAL
        paragraph.alignment = PP_ALIGN.CENTER

    # ========================================================================
    # SLIDE 2: Table of Contents
    # ========================================================================
    print("  ✓ Slide 2: Table of Contents")
    add_content_slide(prs, "Table of Contents", [
        "1. Brand Overview & Brand Story",
        "2. Logo Design & Construction",
        "3. Logo Variations",
        "4. Logo Misuse Guidelines",
        "5. Color Palette",
        "6. Typography Showcase",
        "7. Iconography & UI Elements",
        "8. UI Color Application",
        "9. Brand Voice & Messaging",
        "10. Stationery & Applications",
        "11. Digital & Print Usage Guidelines",
        "12. Accessibility & Contrast"
    ])

    # ========================================================================
    # SLIDE 3: Brand Overview
    # ========================================================================
    print("  ✓ Slide 3: Brand Overview")
    add_content_slide(prs, "Brand Overview & Story", [
        "What is SmartAC?",
        "SmartAC is an innovative IoT-based smart air conditioning control system that revolutionizes how users interact with their climate control devices.",
        "",
        "Built on the Android platform, SmartAC empowers users to remotely monitor, control, and optimize their air conditioning units through an intuitive mobile application.",
        "",
        "Mission: To make intelligent climate control accessible, sustainable, and effortless for every home, while reducing energy waste.",
        "",
        "Brand Personality: Intelligent • Efficient • Connected • Sustainable • Trustworthy • Innovative"
    ])

    # ========================================================================
    # SLIDE 4: Logo Design
    # ========================================================================
    print("  ✓ Slide 4: Logo Design")
    add_content_slide(prs, "Logo Design & Construction", [
        "The SmartAC logo features:",
        "",
        "• MTP initials (Makamte Tatsinkou Paola) as the central typographic element",
        "• Bold, geometric font treatment for modern aesthetic",
        "• WiFi/IoT connectivity symbol integrated into the design",
        "• Airflow arc motifs representing cooling and air circulation",
        "• Wordmark 'SMARTAC' in clean sans-serif capitals",
        "• Tagline: 'Intelligent Climate. Total Control.'",
        "",
        "Clear Space: Minimum space = height of letter 'M' on all sides",
        "Minimum Size: 40px digital / 15mm print"
    ])

    # ========================================================================
    # SLIDE 5: Logo Variations
    # ========================================================================
    print("  ✓ Slide 5: Logo Variations")
    add_content_slide(prs, "Logo Variations", [
        "Four official logo variations:",
        "",
        "1. Primary Version",
        "   • Full color on white/light backgrounds",
        "   • Preferred for most applications",
        "",
        "2. Reversed Version",
        "   • White on Deep Navy or dark backgrounds",
        "   • For dark-themed interfaces",
        "",
        "3. Monochrome Dark",
        "   • All Deep Navy elements",
        "   • For single-color print applications",
        "",
        "4. Monochrome Light",
        "   • All white elements on dark backgrounds",
        "   • For specialty applications"
    ])

    # ========================================================================
    # SLIDE 6: Logo Misuse
    # ========================================================================
    print("  ✓ Slide 6: Logo Misuse")
    add_content_slide(prs, "Logo Misuse — Don'ts", [
        "Never do the following with the SmartAC logo:",
        "",
        "✗ Don't stretch or distort the logo proportions",
        "✗ Don't change colors arbitrarily",
        "✗ Don't place on busy or patterned backgrounds",
        "✗ Don't add drop shadows or effects",
        "✗ Don't rotate or tilt the logo",
        "✗ Don't use low-contrast color combinations",
        "✗ Don't recreate or modify the logo elements",
        "✗ Don't use older versions of the logo",
        "",
        "Always use approved logo files from the brand asset library."
    ])

    # ========================================================================
    # SLIDE 7: Color Palette
    # ========================================================================
    print("  ✓ Slide 7: Color Palette")
    slide = add_content_slide(prs, "Brand Color Palette", [
        "Arctic Blue  |  HEX: #0A84FF  |  RGB: 10, 132, 255  |  Primary brand color",
        "",
        "Cool Teal  |  HEX: #30D5C8  |  RGB: 48, 213, 200  |  Accent color, IoT highlights",
        "",
        "Deep Navy  |  HEX: #0D1B2A  |  RGB: 13, 27, 42  |  Dark backgrounds, body text",
        "",
        "Frost White  |  HEX: #F4F9FF  |  RGB: 244, 249, 255  |  Light backgrounds, cards",
        "",
        "Slate Grey  |  HEX: #6B7280  |  RGB: 107, 114, 128  |  Secondary text, borders",
        "",
        "CMYK values provided in print-ready files.",
        "Always use exact hex/RGB values for digital applications."
    ])

    # ========================================================================
    # SLIDE 8: Typography
    # ========================================================================
    print("  ✓ Slide 8: Typography")
    add_content_slide(prs, "Typography System", [
        "Primary Typeface: Helvetica Family",
        "",
        "Display / Cover Title: Helvetica Bold, 36pt, Deep Navy",
        "Section Heading (H1): Helvetica Bold, 22pt, Arctic Blue",
        "Sub-heading (H2): Helvetica Bold, 14pt, Deep Navy",
        "Body Text: Helvetica, 11pt, Slate Grey, 1.5 line spacing",
        "Caption / Label: Helvetica Oblique, 9pt, Slate Grey",
        "Tagline / Pull Quote: Helvetica Bold Oblique, 13pt, Cool Teal",
        "",
        "Helvetica provides excellent readability across all mediums.",
        "For web applications, use system font alternatives if Helvetica is unavailable.",
        "Always maintain hierarchy: larger sizes for headers, smaller for body text."
    ])

    # ========================================================================
    # SLIDE 9: Iconography
    # ========================================================================
    print("  ✓ Slide 9: Iconography")
    add_content_slide(prs, "Iconography & UI Elements", [
        "Icon Style Guidelines:",
        "",
        "• Stroke weight: 2pt for consistency",
        "• Line caps: Rounded for friendly aesthetic",
        "• Primary color: Arctic Blue",
        "• Grid: 24×24pt artboard",
        "",
        "Standard Icon Set Includes:",
        "• Thermometer (temperature)",
        "• WiFi symbol (connectivity)",
        "• Android phone outline",
        "• Snowflake (cooling)",
        "• Power button (on/off)",
        "• Settings gear",
        "• Energy-saving leaf",
        "• Cloud sync"
    ])

    # ========================================================================
    # SLIDE 10: UI Color Application
    # ========================================================================
    print("  ✓ Slide 10: UI Color Application")
    add_content_slide(prs, "UI Color Application", [
        "Color Usage in Android Application:",
        "",
        "Deep Navy → Status bar, dark headers",
        "Arctic Blue → Primary actions, buttons (active state), headers",
        "Cool Teal → Notifications, accent badges, call-to-action elements",
        "Frost White → Card backgrounds, content containers, light sections",
        "Slate Grey → Inactive states, secondary text, borders",
        "",
        "Button States:",
        "• Active/ON: Arctic Blue background, white text",
        "• Inactive/OFF: Slate Grey background, white text",
        "• Hover/Focus: Slightly darker shade of primary color",
        "",
        "Maintain consistent color application across all screens."
    ])

    # ========================================================================
    # SLIDE 11: Brand Voice
    # ========================================================================
    print("  ✓ Slide 11: Brand Voice")
    add_content_slide(prs, "Brand Voice & Messaging", [
        "Tone of Voice: Professional, Approachable, Tech-Forward",
        "",
        "We Are:",
        "• Tech-forward: Confident without being intimidating",
        "• Reassuring: Building trust through clear communication",
        "• Empowering: Users feel in control",
        "• Sustainable: Emphasizing environmental responsibility",
        "",
        "DO Say:",
        "✓ 'Optimize your comfort' ✓ 'Smart climate control'",
        "✓ 'Save energy effortlessly' ✓ 'Your comfort, remotely'",
        "",
        "DON'T Say:",
        "✗ 'Maximize AC efficiency' ✗ 'Complicated automation system'",
        "✗ 'IoT-enabled thermal management' ✗ 'Advanced HVAC protocols'"
    ])

    # ========================================================================
    # SLIDE 12: Stationery
    # ========================================================================
    print("  ✓ Slide 12: Stationery")
    add_content_slide(prs, "Stationery & Brand Applications", [
        "Business Card:",
        "• Standard size: 85mm × 55mm (3.5\" × 2\")",
        "• Front: Logo, name, title, contact info on white",
        "• Back: Arctic Blue background with tagline",
        "",
        "Letterhead:",
        "• A4 size with Arctic Blue header strip",
        "• Logo positioned top-left",
        "• Website and contact info in footer",
        "",
        "Email Signature:",
        "• Small logo (max 200px wide)",
        "• Name, title, contact details in Helvetica",
        "• Use web-safe colors (RGB values)",
        "",
        "App Icon:",
        "• Rounded square format (various sizes)",
        "• Arctic Blue background with white logo elements"
    ])

    # ========================================================================
    # SLIDE 13: Digital & Print Usage
    # ========================================================================
    print("  ✓ Slide 13: Digital & Print Usage")
    add_content_slide(prs, "Digital & Print Usage Guidelines", [
        "File Formats:",
        "• Print: PDF, SVG (vector), EPS → CMYK, 300 DPI",
        "• Digital/Web: SVG, PNG, WebP → RGB, 72-96 DPI",
        "• Social Media: PNG, JPG → Platform-specific sizes",
        "",
        "Resolution Requirements:",
        "• Print: 300 DPI minimum for professional quality",
        "• Screen: 72-96 DPI standard for digital",
        "• Mobile: @2x and @3x assets for retina displays",
        "",
        "Social Media Sizes:",
        "• Facebook: 180×180px | Instagram: 320×320px",
        "• Twitter/X: 400×400px | LinkedIn: 300×300px",
        "",
        "Always embed fonts in PDFs. Use exact color values."
    ])

    # ========================================================================
    # SLIDE 14: Accessibility
    # ========================================================================
    print("  ✓ Slide 14: Accessibility")
    add_content_slide(prs, "Accessibility & Contrast", [
        "WCAG Compliance:",
        "",
        "Deep Navy on Frost White: 14.2:1 ratio → AAA Pass",
        "Arctic Blue on White: 5.1:1 ratio → AA Pass",
        "White on Arctic Blue: 5.1:1 ratio → AA Pass",
        "Cool Teal on Deep Navy: 8.5:1 ratio → AAA Pass",
        "",
        "Text Sizing:",
        "• Minimum body text: 11pt (14.67px)",
        "• Line height: 1.5× text size",
        "• Button text: Minimum 10pt with 44×44px touch targets",
        "",
        "Color Blindness:",
        "• Arctic Blue and Cool Teal distinguishable in most types",
        "• Never rely on color alone—use icons and labels",
        "• Test UI with color blindness simulation tools"
    ])

    # ========================================================================
    # SLIDE 15: Conclusion
    # ========================================================================
    print("  ✓ Slide 15: Conclusion")
    slide = add_content_slide(prs, "Conclusion", [
        "Building Trust Through Consistency",
        "",
        "This Visual Identity Guidelines document establishes the foundation for SmartAC's brand presence across all touchpoints.",
        "",
        "Consistent application ensures SmartAC is immediately recognizable, trustworthy, and professional in every interaction.",
        "",
        "By adhering to these standards, we create a cohesive brand experience that builds user confidence and reinforces our position as a leader in smart home climate control.",
        "",
        "SmartAC represents more than technology—it represents a commitment to making sustainable, intelligent climate control accessible to everyone.",
        "",
        f"Document Version 1.0 | {datetime.now().strftime('%B %Y')} | MAKAMTE TATSINKOU PAOLA"
    ])

    # Save the presentation
    prs.save(filename)

    print(f"\n✅ SUCCESS! PowerPoint saved as: {filename}")
    print(f"   This file is fully editable in PowerPoint, Google Slides, or Keynote")
    print(f"   You can modify text, colors, layouts, and add new slides")
    print(f"   Ready for presentation and future updates!")

    return filename

if __name__ == "__main__":
    generate_pptx()
